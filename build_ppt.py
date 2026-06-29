"""从 Markdown 简历生成 A4 可编辑 PPT — 与 HTML 简历样式一致"""
import re, sys, io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
md = Path('data/resume_optimized.md').read_text('utf-8')

def get_section(title):
    pattern = rf'## {title}\n(.*?)(?=\n## |\n---|\Z)'
    m = re.search(pattern, md, re.DOTALL)
    return m.group(1).strip() if m else ''

summary = get_section('个人简介')
skills = []
for line in get_section('技能栈').split('\n'):
    m = re.match(r'- \*\*(.+?)\*\*: (.+)', line)
    if m: skills.append((m.group(1), m.group(2)))

exp_text = get_section('工作经历')
exps = []
for block in re.split(r'\n### ', '\n' + exp_text):
    if not block.strip(): continue
    lines = block.strip().split('\n')
    h = lines[0].replace('### ','')
    parts = h.split('|')
    title = parts[0].strip()
    company = parts[1].strip() if len(parts)>1 else ''
    dur = ''; pts = []
    for l in lines[1:]:
        s = l.strip()
        if s.startswith('*'): dur = s.replace('*','').strip()
        elif s.startswith('- ') and '技术栈' not in s: pts.append(s[2:])
    exps.append({'title':title,'company':company,'duration':dur,'points':pts})

proj_text = get_section('项目经验')
projs = []
for block in re.split(r'\n### ', '\n' + proj_text):
    if not block.strip(): continue
    lines = block.strip().split('\n')
    name = lines[0].replace('### ','')
    role = url = ''; pts = []
    for l in lines[1:]:
        s = l.strip()
        if s.startswith('*'): role = s.replace('*','').strip()
        elif s.startswith('链接:'): url = s.split(':',1)[1].strip()
        elif s.startswith('- ') and '技术栈' not in s: pts.append(s[2:])
    projs.append({'name':name,'role':role,'url':url,'points':pts})

edu_text = get_section('教育背景')
edu_lines = edu_text.split('\n')
edu_school = edu_lines[0].replace('**','').split('|')[0].strip()
edu_degree = edu_lines[0].split('|')[1].strip() if len(edu_lines[0].split('|'))>1 else ''
edu_major = edu_lines[0].split('|')[2].strip() if len(edu_lines[0].split('|'))>2 else ''

job_lines = [l.strip('- ').strip() for l in get_section('求职意向').split('\n') if l.startswith('-')]
job_dict = {}
for l in job_lines:
    if '：' in l or ':' in l:
        sep = '：' if '：' in l else ':'; k,v = l.split(sep,1); job_dict[k.strip()] = v.strip()

portfolio_text = get_section('作品集')
port_lines = [l.strip('- ').strip() for l in portfolio_text.split('\n') if l.strip().startswith('-')]
port_line = ' | '.join(port_lines)

target_match = re.search(r'目标方向[：:]\s*(.+)', md)
target = target_match.group(1).strip() if target_match else 'AIGC内容创作者'

# ━━━━ Colors ━━━━
DARK_BLUE1 = RGBColor(0x1a, 0x1a, 0x2e)
DARK_BLUE2 = RGBColor(0x16, 0x21, 0x3e)
DEEP_BLUE  = RGBColor(0x0f, 0x34, 0x60)
ACCENT_RED = RGBColor(0xe9, 0x45, 0x60)
GREEN      = RGBColor(0x5a, 0x9a, 0x5a)
DARK       = RGBColor(0x1a, 0x1a, 0x1a)
GRAY       = RGBColor(0x55, 0x55, 0x55)
LGRAY      = RGBColor(0x88, 0x88, 0x88)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG   = RGBColor(0xf8, 0xf9, 0xfa)

# ━━━━ A4 Setup ━━━━
prs = Presentation()
prs.slide_width  = Cm(21.0)   # A4 width
prs.slide_height = Cm(29.7)   # A4 height
W = prs.slide_width
H = prs.slide_height
MARGIN_LEFT = Cm(1.5)
MARGIN_RIGHT = Cm(1.5)
CONTENT_W = W - MARGIN_LEFT - MARGIN_RIGHT

def add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, left, top, width, height, text, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.alignment = align
    return tf

def add_multiline(slide, left, top, width, height, lines, size=12, color=DARK, spacing=Pt(4)):
    """Add text with multiple paragraphs"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = spacing
    return tf

def estimated_height(lines, size=12, spacing=4, line_height=1.4):
    """Estimate text height in cm"""
    cm_per_pt = 0.0353
    return len(lines) * (size * line_height + spacing) * cm_per_pt / 10

# ━━━━ Build slides ━━━━
current_slide = None
current_y = Cm(0)

def new_slide():
    global current_slide, current_y
    current_slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = current_slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    current_y = Cm(0)
    return current_slide

def draw_header(title_text, subtitle=''):
    global current_y, current_slide
    # Blue gradient header bar
    bar_h = Cm(3.0)
    bar = add_shape(current_slide, Cm(0), Cm(0), W, bar_h, DARK_BLUE1)
    add_text(current_slide, MARGIN_LEFT, Cm(0.6), CONTENT_W, Cm(1.0), title_text, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(current_slide, MARGIN_LEFT, Cm(1.6), CONTENT_W, Cm(0.6), subtitle, size=12, color=RGBColor(0xcc,0xcc,0xdd))
    # Photo placeholder
    photo = add_shape(current_slide, W - Cm(4.5), Cm(0.6), Cm(2.8), Cm(3.8), RGBColor(0x25, 0x25, 0x40))
    add_text(current_slide, W - Cm(4.5), Cm(2.0), Cm(2.8), Cm(0.5), '证件照', size=10, color=RGBColor(0x88,0x88,0x99), align=PP_ALIGN.CENTER)
    current_y = Cm(3.5)

def draw_section_title(title):
    global current_y, current_slide
    # Green left border
    add_shape(current_slide, MARGIN_LEFT, current_y, Cm(0.15), Cm(0.7), GREEN)
    add_text(current_slide, MARGIN_LEFT + Cm(0.4), current_y + Cm(0.05), CONTENT_W, Cm(0.6), title, size=15, bold=True)
    current_y += Cm(0.9)

def check_space(needed_cm):
    """If not enough space on current slide, make new one"""
    global current_y, current_slide
    if current_y + Cm(needed_cm) > H - Cm(2.0):
        current_slide = new_slide()
        add_text(current_slide, MARGIN_LEFT, Cm(0.4), CONTENT_W, Cm(0.4), '（续）', size=9, color=LGRAY, align=PP_ALIGN.RIGHT)
        current_y = Cm(1.2)
        return True
    return False

# ── Slide 1: Header + Summary + Skills ──
new_slide()
draw_header(target, 'AIGC内容创作 · AI科技自媒体 · 深圳 · 15-20K')

# Summary
draw_section_title('个人概述')
summary_lines = [summary]
add_multiline(current_slide, MARGIN_LEFT, current_y, CONTENT_W, Cm(1.5), summary_lines, size=11, color=GRAY)
# Estimate: summary is about 2-3 lines
current_y += Cm(1.6)

# Skills with tags
draw_section_title('核心技能')
tag_y = current_y
has_skill_overflow = False
for level, items_str in skills:
    items = [s.strip() for s in items_str.split(',')]
    add_text(current_slide, MARGIN_LEFT, current_y, Cm(1.8), Cm(0.5), level, size=11, bold=True)
    x = MARGIN_LEFT + Cm(2.2)
    first_row = True
    for item in items:
        tag_w = Cm(len(item) * 0.42 + 0.8)
        if x + tag_w > W - MARGIN_RIGHT:
            x = MARGIN_LEFT + Cm(2.2)
            current_y += Cm(0.7)
            first_row = False
        add_shape(current_slide, x, current_y + Cm(0.05), tag_w, Cm(0.5), DARK_BLUE1)
        add_text(current_slide, x + Cm(0.1), current_y + Cm(0.06), tag_w - Cm(0.2), Cm(0.45), item, size=9, color=WHITE)
        x += tag_w + Cm(0.2)
    if not first_row:
        current_y += Cm(0.3)
    current_y += Cm(0.9)

current_y += Cm(0.5)

# ── Work Experience ──
check_space(6)
draw_section_title('工作经历')
for exp in exps:
    needed = len(exp['points']) * 0.8 + 1.8
    check_space(needed)
    hdr = f"{exp['title']}  |  {exp['company']}"
    add_text(current_slide, MARGIN_LEFT, current_y, CONTENT_W, Cm(0.5), hdr, size=13, bold=True)
    current_y += Cm(0.5)
    add_text(current_slide, MARGIN_LEFT, current_y, CONTENT_W, Cm(0.4), exp['duration'], size=10, color=GREEN)
    current_y += Cm(0.5)
    bullets = [f"• {p}" for p in exp['points']]
    add_multiline(current_slide, MARGIN_LEFT + Cm(0.3), current_y, CONTENT_W - Cm(0.3), Cm(needed), bullets, size=10, color=DARK, spacing=Pt(3))
    current_y += Cm(needed)

# ── Projects ──
check_space(4)
draw_section_title('项目经验')
for proj in projs:
    needed = len(proj['points']) * 0.7 + 1.5
    check_space(needed)
    hdr = proj['name']
    if proj['role']: hdr += f"  —  {proj['role']}"
    add_text(current_slide, MARGIN_LEFT, current_y, CONTENT_W, Cm(0.5), hdr, size=13, bold=True)
    current_y += Cm(0.5)
    bullets = [f"• {p}" for p in proj['points']]
    add_multiline(current_slide, MARGIN_LEFT + Cm(0.3), current_y, CONTENT_W - Cm(0.3), Cm(needed), bullets, size=10, color=DARK, spacing=Pt(3))
    current_y += Cm(needed)

# ── Education + Portfolio + Job Target ──
check_space(5)
draw_section_title('教育背景')
add_text(current_slide, MARGIN_LEFT, current_y, CONTENT_W, Cm(0.5), f'{edu_school}  |  {edu_degree}  |  {edu_major}', size=12, bold=True)
current_y += Cm(0.8)

draw_section_title('作品集 & 社区')
add_text(current_slide, MARGIN_LEFT, current_y, CONTENT_W, Cm(0.6), port_line, size=10, color=GRAY)
current_y += Cm(1.0)

draw_section_title('求职意向')
job_parts = [f"{k}: {v}" for k, v in job_dict.items()]
add_multiline(current_slide, MARGIN_LEFT, current_y, CONTENT_W, Cm(1.5), job_parts, size=11, color=DARK)
current_y += Cm(1.5)

# Footer
add_text(current_slide, MARGIN_LEFT, H - Cm(1.2), CONTENT_W, Cm(0.5), '叁十三 · AIGC创作者  |  深圳  |  sanshisanAIGC', size=9, color=LGRAY, align=PP_ALIGN.CENTER)

# ━━━━ Save ━━━━
output = Path('data/resume.pptx')
prs.save(str(output))
print(f'PPT 已生成: {output}')
print(f'A4 尺寸 | 共 {len(prs.slides)} 页 | 与HTML简历样式一致')
